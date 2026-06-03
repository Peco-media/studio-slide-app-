"""
slide_generator.py
スタジオ情報 → pptxスライドを生成するモジュール
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from PIL import Image
import io
import copy


# ── カラー定義 ──────────────────────────────
NAVY   = RGBColor(0x2C, 0x4A, 0x6E)
BLUE   = RGBColor(0x5B, 0x8D, 0xB8)
BLUE_L = RGBColor(0xD6, 0xE8, 0xF5)
BLUE_XL= RGBColor(0xEE, 0xF6, 0xFC)
YELLOW = RGBColor(0xE8, 0xB8, 0x4B)
YELLOW_L=RGBColor(0xFD, 0xF6, 0xE3)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GRAY   = RGBColor(0x6B, 0x72, 0x80)
DARK   = RGBColor(0x1E, 0x29, 0x3B)
TEAL   = RGBColor(0x2A, 0x7A, 0x6F)
TEAL_L = RGBColor(0xD1, 0xF0, 0xEC)
RED    = RGBColor(0xC0, 0x39, 0x2B)
RED_L  = RGBColor(0xFD, 0xEC, 0xEA)


def inch(val):
    return Inches(val)

def pt(val):
    return Pt(val)


def add_rect(slide, x, y, w, h, fill_color, line_color=None, line_width=0):
    shape = slide.shapes.add_shape(1, inch(x), inch(y), inch(w), inch(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width) if line_width else Pt(0.5)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h, font_size, bold=False,
             color=None, align=PP_ALIGN.LEFT, valign=None):
    txBox = slide.shapes.add_textbox(inch(x), inch(y), inch(w), inch(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = pt(font_size)
    run.font.bold = bold
    run.font.name = "Arial"
    if color:
        run.font.color.rgb = color
    return txBox


def add_image_from_bytes(slide, img_bytes, x, y, w, h):
    """PIL Imageまたはbytesからスライドに画像を追加（coverクロップ）"""
    if isinstance(img_bytes, bytes):
        img = Image.open(io.BytesIO(img_bytes))
    else:
        img = img_bytes

    # coverリサイズ: アスペクト比を保ちつつ指定サイズに合わせてクロップ
    target_w = int(w * 96)   # 96dpi相当
    target_h = int(h * 96)
    orig_w, orig_h = img.size
    scale = max(target_w / orig_w, target_h / orig_h)
    new_w = int(orig_w * scale)
    new_h = int(orig_h * scale)
    img = img.resize((new_w, new_h), Image.LANCZOS)
    left = (new_w - target_w) // 2
    top  = (new_h - target_h) // 2
    img = img.crop((left, top, left + target_w, top + target_h))

    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    slide.shapes.add_picture(buf, inch(x), inch(y), inch(w), inch(h))


def add_tag(slide, text, x, y, color=BLUE, bg=BLUE_XL):
    w = min(len(text) * 0.155 + 0.35, 3.2)
    add_rect(slide, x, y, w, 0.28, bg, color)
    add_text(slide, text, x, y, w, 0.28, 8.5, color=NAVY, align=PP_ALIGN.CENTER)
    return w


def generate_slides(studios: list, client_name: str = "") -> bytes:
    """
    studios: [
      {
        "name": str,
        "area": str,
        "address": str,
        "access": str,
        "size": str,
        "ceiling": str,
        "rooms": str,
        "kitchen": str,
        "natural_light": str,
        "electric": str,
        "wifi": str,
        "min_time": str,
        "price_still": str,
        "price_movie": str,
        "tags": [str, ...],
        "note": str,
        "caution": str,
        "url": str,
        "photos": [PIL.Image, ...],   ← アップロード画像（リスト）
      },
      ...
    ]
    戻り値: pptxのbytes
    """
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(5.625)

    blank_layout = prs.slide_layouts[6]  # 完全ブランク

    # ━━━ スライド1: タイトル ━━━
    sl = prs.slides.add_slide(blank_layout)
    add_rect(sl, 0, 0, 10, 5.625, BLUE_XL)
    add_rect(sl, 0, 0, 10, 0.55, BLUE)
    header = client_name + "　｜　" if client_name else ""
    add_text(sl, header + "ペット可 撮影スタジオ候補",
             0.4, 0.08, 9, 0.38, 13, bold=True, color=WHITE)
    add_text(sl, "スタジオ候補　" + str(len(studios)) + "件",
             0.6, 0.85, 8, 0.9, 34, bold=True, color=NAVY)
    add_rect(sl, 0.6, 1.82, 5, 0.04, YELLOW)

    conditions = ["🐾 ペット可", "📍 東京都内", "☀️ 自然光あり", "💴 〜30,000円"]
    bx = 0.6
    for c in conditions:
        add_rect(sl, bx, 2.05, 1.9, 0.42, WHITE, BLUE_L)
        add_text(sl, c, bx, 2.05, 1.9, 0.42, 12, color=NAVY, align=PP_ALIGN.CENTER)
        bx += 2.05

    # スタジオミニカード（タイトルスライド下部）
    colors = [BLUE, TEAL, RGBColor(0x85, 0x4F, 0x0B), RGBColor(0x53, 0x4A, 0xB7)]
    for i, s in enumerate(studios[:4]):
        col = colors[i % len(colors)]
        cx = 0.4 + i * (9.2 / max(len(studios), 1)) if len(studios) > 2 else 0.4 + i * 4.75
        cw = min(4.45, 9.2 / len(studios) - 0.1)
        add_rect(sl, cx, 2.65, cw, 2.55, WHITE, BLUE_L)
        add_rect(sl, cx, 2.65, cw, 0.42, col)
        add_text(sl, f"{'①②③④'[i]} {s['name']}", cx+0.1, 2.68, cw-0.2, 0.36,
                 12, bold=True, color=WHITE)
        # 写真
        if s.get("photos"):
            try:
                add_image_from_bytes(sl, s["photos"][0], cx+0.1, 3.15, cw-0.2, 1.9)
            except Exception:
                pass

    # ━━━ スライド2: 比較表 ━━━
    sl = prs.slides.add_slide(blank_layout)
    add_rect(sl, 0, 0, 10, 5.625, WHITE)
    add_rect(sl, 0, 0, 10, 0.55, BLUE)
    hdr = (client_name + "　｜　") if client_name else ""
    add_text(sl, hdr + "ペット可 撮影スタジオ候補", 0.4, 0.08, 9.2, 0.38, 12, bold=True, color=WHITE)
    add_text(sl, "スタジオ比較", 0.4, 0.68, 8, 0.5, 20, bold=True, color=NAVY)
    add_rect(sl, 0.4, 1.22, 2.5, 0.04, YELLOW)

    from pptx.util import Inches as I
    from pptx.oxml.ns import qn
    import lxml.etree as etree

    # テーブルをpython-pptxで作成
    row_labels = ["エリア","アクセス","広さ","自然光","ペット可","控室","キッチン","最低時間","スチール料金"]
    row_keys   = ["area","access","size","natural_light","pet","rooms","kitchen","min_time","price_still"]
    pet_key = "pet"

    n_studios = len(studios)
    col_widths = [Inches(2.2)] + [Inches(7.0 / n_studios)] * n_studios
    tbl = sl.shapes.add_table(
        len(row_labels)+1, n_studios+1,
        Inches(0.3), Inches(1.35),
        Inches(9.4), Inches(4.1)
    ).table

    # 列幅
    for ci, cw in enumerate(col_widths):
        tbl.columns[ci].width = cw

    def set_cell(cell, text, bg=None, txt_color=None, bold=False, font_size=10, align=PP_ALIGN.CENTER):
        cell.text = text
        p = cell.text_frame.paragraphs[0]
        p.alignment = align
        run = p.runs[0] if p.runs else p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.name = "Arial"
        if txt_color:
            run.font.color.rgb = txt_color
        if bg:
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg

    # ヘッダー行
    set_cell(tbl.cell(0, 0), "項目", bg=NAVY, txt_color=WHITE, bold=True)
    header_colors = [BLUE, TEAL, RGBColor(0x85, 0x4F, 0x0B), RGBColor(0x53, 0x4A, 0xB7)]
    nums = "①②③④"
    for i, s in enumerate(studios):
        set_cell(tbl.cell(0, i+1), f"{nums[i]} {s['name']}", bg=header_colors[i%len(header_colors)], txt_color=WHITE, bold=True)

    # データ行
    for ri, (label, key) in enumerate(zip(row_labels, row_keys)):
        bg = WHITE if ri % 2 == 0 else BLUE_XL
        set_cell(tbl.cell(ri+1, 0), label, bg=bg, txt_color=NAVY, bold=True, align=PP_ALIGN.LEFT)
        for ci, s in enumerate(studios):
            val = s.get(key, "—")
            col = DARK
            if val and ("✔" in val or "あり" in val or "OK" in val):
                col = TEAL
            elif val and ("⚠" in val or "要" in val or "なし" in val):
                col = RED
            set_cell(tbl.cell(ri+1, ci+1), val or "—", bg=bg, txt_color=col)

    # ━━━ スライド3〜: 各スタジオ詳細 ━━━
    for idx, s in enumerate(studios):
        col = header_colors[idx % len(header_colors)]
        sl = prs.slides.add_slide(blank_layout)
        add_rect(sl, 0, 0, 10, 5.625, WHITE)

        # ヘッダー
        add_rect(sl, 0, 0, 10, 1.0, col)
        num = "①②③④"[idx]
        add_text(sl, f"{num}　{s['name']}", 0.35, 0.06, 8.0, 0.52, 20, bold=True, color=WHITE)
        add_text(sl, s.get("address",""), 0.35, 0.62, 7.5, 0.28, 10, color=RGBColor(0xD9,0xEA,0xF7))
        if s.get("url"):
            add_text(sl, s["url"], 7.1, 0.65, 2.7, 0.24, 8, color=RGBColor(0xD9,0xEA,0xF7))

        # 左: 基本情報
        info = [
            ("📍 アクセス",  s.get("access","")),
            ("📐 広さ",      f"{s.get('size','')}　天井高 {s.get('ceiling','')}"),
            ("🚿 控室",      s.get("rooms","")),
            ("🍳 キッチン",  s.get("kitchen","")),
            ("☀️ 自然光",   s.get("natural_light","")),
            ("🔌 電気",      s.get("electric","")),
            ("📶 Wi-Fi",     s.get("wifi","")),
            ("🕐 最低時間",  s.get("min_time","")),
        ]
        iy = 1.08
        for label, val in info:
            add_text(sl, label, 0.3, iy, 1.6, 0.32, 9.5, bold=True, color=NAVY)
            add_text(sl, val,   1.95, iy, 3.1, 0.32, 9,   color=DARK)
            iy += 0.33

        # 料金ボックス
        add_rect(sl, 0.3, iy+0.04, 4.8, 0.72, BLUE_XL, BLUE_L)
        add_text(sl, "💴 スチール撮影", 0.42, iy+0.08, 2.2, 0.26, 10, bold=True, color=NAVY)
        add_text(sl, s.get("price_still",""), 0.42, iy+0.34, 4.5, 0.36, 9.5, color=DARK)

        # 右: 写真
        photos = s.get("photos", [])
        if photos:
            # メイン写真（大）
            try:
                add_image_from_bytes(sl, photos[0], 5.35, 1.05, 4.45, 2.2)
            except Exception:
                add_rect(sl, 5.35, 1.05, 4.45, 2.2, BLUE_XL, BLUE_L)
                add_text(sl, "写真1", 5.35, 1.05, 4.45, 2.2, 11, color=BLUE, align=PP_ALIGN.CENTER)
            # サムネイル
            for pi in range(min(3, len(photos)-1)):
                try:
                    add_image_from_bytes(sl, photos[pi+1], 5.35+pi*1.52, 3.32, 1.42, 0.88)
                except Exception:
                    add_rect(sl, 5.35+pi*1.52, 3.32, 1.42, 0.88, BLUE_XL, BLUE_L)
        else:
            add_rect(sl, 5.35, 1.05, 4.45, 2.2, BLUE_XL, BLUE_L)
            add_text(sl, "写真をアップロードしてください", 5.35, 1.05, 4.45, 2.2, 10,
                     color=GRAY, align=PP_ALIGN.CENTER)

        # 確認事項
        add_rect(sl, 5.35, 4.28, 4.45, 0.72, YELLOW_L, YELLOW)
        add_text(sl, "📋 確認事項", 5.48, 4.31, 4.1, 0.26, 10, bold=True, color=NAVY)
        add_text(sl, s.get("caution",""), 5.48, 4.57, 4.1, 0.38, 9, color=DARK)

        # 特徴タグ
        add_text(sl, "スタジオの特徴", 0.3, 4.72, 4.5, 0.26, 10, bold=True, color=NAVY)
        tx, ty = 0.3, 5.03
        for tag in s.get("tags", []):
            tw = min(len(tag) * 0.155 + 0.35, 3.2)
            if tx + tw > 5.1:
                tx, ty = 0.3, ty + 0.33
            add_tag(sl, tag, tx, ty, color=col, bg=BLUE_XL)
            tx += tw + 0.1

    # bytesで返す
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()
